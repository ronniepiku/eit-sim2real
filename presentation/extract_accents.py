"""One-off: extract the Designer-generated corner accent graphics from slide 1
(added by hand in PowerPoint) into reusable XML assets, so build_deck.py can
clone them onto every slide programmatically and the design stays reproducible."""
from pathlib import Path
from lxml import etree
from pptx import Presentation

HERE = Path(__file__).parent
ASSETS = HERE / "assets"
ASSETS.mkdir(exist_ok=True)

prs = Presentation(HERE / "EE52037_Piku_Presentation.pptx")
s = prs.slides[0]
names = {"Group 13": "accent_topright.xml", "Group 19": "accent_bottomleft.xml"}
for shp in s.shapes:
    if shp.name in names:
        xml = etree.tostring(shp._element)
        (ASSETS / names[shp.name]).write_bytes(xml)
        print(f"wrote {names[shp.name]}: {len(xml)} bytes, pos=({shp.left},{shp.top}) size=({shp.width},{shp.height})")

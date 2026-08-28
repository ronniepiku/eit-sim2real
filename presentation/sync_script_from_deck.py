"""Regenerate speaker_script.md from the pptx's own (authoritative) notes, so
the two stay in sync after editing notes directly in build_deck.py. Inverse
of apply_script.py's usual direction; run this, then apply_script.py is a
no-op confirmation."""
from pathlib import Path
from pptx import Presentation

HERE = Path(__file__).parent
prs = Presentation(HERE / "EE52037_Piku_Presentation.pptx")

total_words = 0
blocks = []
for s in prs.slides:
    text = s.notes_slide.notes_text_frame.text
    header, body = text.split("]", 1)
    header = header.lstrip("[").strip()
    body = body.strip()
    total_words += len(body.split())
    blocks.append(f"## Slide {len(blocks) + 1} — {header}\n\n{body}\n")

wpm = 135
mins, secs = divmod(round(total_words / wpm * 60), 60)

header = f"""# EE52037 — presentation script
**Ronald Piku · A Physically Motivated Noise Model for Robust EIT Tactile Classification Under Simulated Domain Shift**

Total {total_words} words ≈ **{mins}:{secs:02d} at 135 words/min** (range: {total_words/145:.1f}–{total_words/125:.1f} min across 125–145 wpm). Slide 11 (References) is not narrated — let it sit on screen during the closing line.

Markers show elapsed time at the *start* of each slide, at 135 wpm.

---

"""
body_md = "\n---\n\n".join(blocks)
footer = """
---

## Recording notes

- **Slides 7 and 8 carry the marks.** Critical Thinking and Argumentation is
  70% of this assessment. Slide 7 is the protocol-inversion finding (the
  main contribution) and slide 8 decomposes it further with the confusion
  matrix. Do not rush these to save time elsewhere.
- **Rehearsal note on research reflexivity.** The self-testing content that
  used to be its own slide (the permutation test and the ρ-sweep,
  confirming the CNN's advantage is really about measurement structure, and
  that the electrode-bias result isn't inflated by cancellation) is no
  longer in the pre-recorded deck. Both checks are still in the
  dissertation (Section on testing assumptions) and are strong, concrete
  material for the 13-minute viva if a question about evidence or
  robustness comes up — keep them ready to describe verbally.
- **Numbers to land clearly:** 94 → 20 → 76 (slide 6); the matched vs
  deployment inversion (slide 7); the no-contact recovery in the confusion
  matrix (slide 8).
- **Slide 11 (References)** needs no delivery time — say the thank-you/
  questions line while it's on screen.
"""

md = header + body_md + footer
(HERE / "speaker_script.md").write_text(md, encoding="utf-8")
print(f"wrote speaker_script.md: {total_words} words, ~{mins}:{secs:02d} at 135 wpm")

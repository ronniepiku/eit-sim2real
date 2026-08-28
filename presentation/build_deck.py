"""Build the EE52037 7-minute viva presentation.

Deck is deliberately sparse: one claim per slide, a single visual carrying the
evidence, and the argument delivered in the voiceover rather than read off the
screen. Speaker notes carry the timed script (kept in sync with
speaker_script.md via apply_script.py, which must be run after this).

Design: a set of corner "accent" motifs (freeform gradient swooshes, matching
the University of Bath colour palette) were created once in PowerPoint
Designer on the title slide and extracted to assets/*.xml by
extract_accents.py. This script clones them onto every slide at a small,
consistent scale so the chrome is uniform and reproducible from code.
"""
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.shapes.group import GroupShape

HERE = Path(__file__).parent
FIG = HERE / "figures"
ASSETS = HERE / "assets"
REC = HERE.parent / "results" / "dataset_validation" / "reconstructions"

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
RULE = RGBColor(0xE6, 0xE6, 0xE3)
BOXFILL = RGBColor(0xF2, 0xF4, 0xF7)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

SLIDE_COUNT = 11  # kept in sync manually; used nowhere critical, just FYI


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = SURFACE
    return s


def tb(s, x, y, w, h, text, size=18, bold=False, italic=False, color=INK,
       align=PP_ALIGN.LEFT, space_after=6, line=1.15, font="Calibri"):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
            r.font.name = font
    return box


def heading(s, title, kicker=None, title_h=0.58):
    if kicker:
        tb(s, 0.85, 0.52, 10.6, 0.40, kicker.upper(), size=13, bold=True, color=BLUE)
        tb(s, 0.85, 0.92, 11.6, title_h, title, size=32, bold=True)
    else:
        tb(s, 0.85, 0.72, 11.6, title_h, title, size=32, bold=True)


def rule(s, y=1.92):
    ln = s.shapes.add_shape(1, Inches(0.85), Inches(y), Inches(11.6), Pt(1.6))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RULE
    ln.line.fill.background()
    ln.shadow.inherit = False


def pic(s, path, x, y, w, crop_top=0.0):
    p = s.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    if crop_top:
        p.crop_top = crop_top
    return p


def caption(s, x, y, w, text, align=PP_ALIGN.LEFT):
    return tb(s, x, y, w, 0.35, text, size=12.5, italic=True, color=MUTED,
              align=align, space_after=0, line=1.2)


def box(s, x, y, w, h, fill=BOXFILL):
    r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def note(s, text):
    s.notes_slide.notes_text_frame.text = text


def num(s, n):
    tb(s, 12.35, 6.92, 0.6, 0.3, str(n), size=11, color=MUTED, align=PP_ALIGN.RIGHT)


_ACCENT_CACHE = {}


def _load_accent(name):
    if name not in _ACCENT_CACHE:
        _ACCENT_CACHE[name] = (ASSETS / name).read_bytes()
    return _ACCENT_CACHE[name]


def accent(s, name, left, top, width, height):
    """Clone a Designer-generated corner accent onto a slide at a given box."""
    grp = parse_xml(_load_accent(name))
    s.shapes._spTree.append(grp)
    gs = GroupShape(grp, s.shapes)
    gs.left, gs.top, gs.width, gs.height = (
        Inches(left), Inches(top), Inches(width), Inches(height))
    return gs


def chrome(s):
    """Standard subtle corner accents used on every content slide (2-11)."""
    accent(s, "accent_topright.xml", 11.55, -0.10, 1.85, 1.13)
    accent(s, "accent_bottomleft.xml", -0.15, 6.62, 1.30, 0.98)


# ═══════════════════════════════════════════════════════ 1 — title
s = slide()
tb(s, 0.85, 0.64, 10.75, 0.50, "MSc Artificial Intelligence in Engineering and Design – EE52037",
   size=15, bold=True, color=BLUE)
tb(s, 0.85, 2.44, 11.20, 2.27,
   "A Physically Motivated Noise Model for\nRobust EIT Tactile Classification Under Simulated Domain Shift",
   size=40, bold=True, line=1.12)
tb(s, 0.85, 5.65, 11.20, 1.19,
   "Author: Ronald Piku       \nSupervisor: Prof. Manuchehr Soleimani\n"
   "Department of Electronic and Electrical Engineering, University of Bath",
   size=14, color=MUTED, line=1.5)
accent(s, "accent_topright.xml", 9.07 / 1, 0.0, 4.27, 2.61)
accent(s, "accent_bottomleft.xml", 0.0, 5.12, 3.17, 2.38)
note(s, """[Title – 0:00 – ~12s]

Hi, I'm Ronald Piku. Thank you for joining me as I present my dissertation: a physically motivated noise model for robust EIT tactile classification under simulated domain shift.""")

# ═══════════════════════════════════════════════════════ 2 — introduction / roadmap
s = slide()
chrome(s)
heading(s, "What This Presentation Covers", kicker="Introduction")
rule(s)
tb(s, 0.85, 2.20, 11.60, 0.85,
   "This project asks whether a physically motivated noise model can close the "
   "simulation-to-reality gap in EIT tactile sensing, and which noise sources "
   "actually need to be in it.",
   size=19, line=1.35)

roadmap = [
    ("1", "The problem and the gap",
     "why simulated EIT classifiers fail on realistic measurements, and what “the gap” means."),
    ("2", "Three results",
     "what closes the gap, how it was tested, and the minimal model that is left standing."),
    ("3", "Testing the claims, and their limits",
     "stress-testing assumptions, then reviewing the limitations."),
]
y = 3.35
for n, t, d in roadmap:
    tb(s, 0.85, y, 0.55, 0.55, n, size=24, bold=True, color=BLUE)
    tb(s, 1.55, y + 0.02, 10.8, 0.40, t, size=19, bold=True)
    tb(s, 1.55, y + 0.46, 10.8, 0.55, d, size=15.5, color=MUTED, line=1.25)
    y += 1.10
num(s, 2)
note(s, """[Introduction – 0:12 – ~23s]

Here's how this presentation is structured. First, I define the problem this
creates and what I mean by the simulation-to-reality gap. Then three results
follow: what closes that gap, how it was evaluated, and the smallest noise
model that still works. Finally, I review the conclusions and discuss the
work's limitations.""")

# ═══════════════════════════════════════════════════════ 3 — the problem (+ diagram)
s = slide()
chrome(s)
heading(s, "How do we get Prosthetics to Feel?", kicker="The problem")
rule(s)
tb(s, 0.85, 2.35, 5.30, 3.55,
   "57.7 million people live with limb loss (McDonald et al., 2021).\n\n"
   "More than a third abandon myoelectric\nprostheses — the absence of touch\n"
   "is a leading reason (Biddiss and Chau, 2007).\n\n"
   "A prosthetic arm has no way of sensing how it is touched or handled, by its "
   "wearer or by someone else, without looking at it.",
   size=17.5, line=1.32, space_after=2)
tb(s, 6.90, 2.28, 5.60, 0.62,
   "Electrical Impedance Tomography (EIT) turns a sheet of soft material "
   "into one continuous touch sensor:",
   size=15.5, color=MUTED, line=1.22)
pic(s, FIG / "p_eit_diagram.png", 6.90, 3.00, 5.60)
caption(s, 6.90, 6.35, 5.60, "Figure 1. EIT sensing principle applied to a prosthetic arm.")
num(s, 3)
note(s, """[The problem – 0:35 – ~56s]

Around fifty-eight million people live with limb loss, and more than a third
abandon a powered prosthetic arm — the most cited reason is that the device
itself has no sense of touch, so it can't tell when or how it is being
handled.

Electrical Impedance Tomography, or EIT, gives it that awareness. As Figure 1
shows, sixteen electrodes sit around a soft sensing disc covering the arm,
modelled on an ionic hydrogel skin (Costa Cornellà et al., 2023) for its
sensing fidelity, even though it isn't the most durable option available.
Inject a small current between one pair, measure the voltage at another —
where somebody presses, that voltage shifts, and the full set of
measurements is a vector you can feed straight into a classifier.""")

# ═══════════════════════════════════════════════════════ 4 — where the field is (+ gap definition)
s = slide()
chrome(s)
heading(s, "Excellent simulation performance, Unreliable in Deployment", kicker="Where the field is",
        title_h=0.90)
rule(s)
tb(s, 0.85, 2.20, 11.60, 0.60,
   "Published EIT touch classifiers report 98–99% accuracy (Park et al., 2022; "
   "Chen et al., 2023). Almost all are trained on simulated measurements, because "
   "labelled real data is expensive.",
   size=17, color=MUTED, line=1.28)

box(s, 0.85, 2.95, 11.50, 1.10)
tb(s, 1.25, 3.06, 10.80, 0.90,
   "The sim2real gap: the drop in performance when a model trained on idealised "
   "simulated data meets the imperfections of real hardware — well documented in "
   "other ML domains such as robotics and computer vision (Tobin et al., 2017), "
   "but rarely modelled in EIT tactile sensing.",
   size=15.5, bold=True, line=1.28)

tb(s, 0.85, 4.28, 11.60, 0.40,
   "Within EIT, four sources dominate:", size=17, bold=True)
for i, (t, d) in enumerate([
        ("Thermal noise", "in the amplifiers"),
        ("Contact impedance", "varies electrode to electrode"),
        ("Electrode placement", "is never exactly where the model assumes"),
        ("Quantisation", "from the analogue-to-digital converter")]):
    y = 4.78 + i * 0.47
    tb(s, 1.15, y, 3.3, 0.40, t, size=15.5, bold=True, color=BLUE)
    tb(s, 4.5, y, 8.0, 0.40, d, size=15.5, color=MUTED)
tb(s, 0.85, 6.63, 11.60, 0.35,
   "Prior work augments predominantly with Gaussian noise only, if any.",
   size=14.5, bold=True)
num(s, 4)
note(s, """[Where the field is – 1:32 – ~47s]

Published EIT touch classifiers report ninety-eight, ninety-nine per cent
accuracy, almost all trained on simulated measurements, because labelled real
data is expensive to collect.

But a simulator gives you a perfect measurement, and hardware never does —
the same gap Tobin and colleagues closed in robotics and computer vision by
training on deliberately varied simulation, and the gap this research aims to
close for EIT. Within EIT, four sources dominate: thermal noise, contact
impedance, electrode placement, and quantisation.

When prior work adds noise at all, it's almost always just Gaussian noise.
This project asks whether that's enough, and if not, which of the four
actually matter.""")

# ═══════════════════════════════════════════════════════ 5 — the gap, made visible
s = slide()
chrome(s)
heading(s, "What Noise Does to the Signal", kicker="Visualising the Gap")
rule(s)
tb(s, 1.35, 2.15, 4.60, 0.40, "Clean simulation", size=18, bold=True, color=BLUE,
   align=PP_ALIGN.CENTER)
tb(s, 7.35, 2.15, 4.60, 0.33, "With noise", size=18, bold=True, color=ORANGE,
   align=PP_ALIGN.CENTER)
pic(s, REC / "random_reconstructed_class_images_clean.png", 0.90, 2.65, 5.50, crop_top=0.06)
pic(s, REC / "random_reconstructed_class_images_noisy.png", 6.90, 2.65, 5.50, crop_top=0.06)
tb(s, 0.85, 6.12, 11.60, 0.40,
   "Same five touches. The contact is plainly visible on the left and gone on the right.",
   size=16.5, color=MUTED)
caption(s, 0.85, 6.55, 11.60,
        "Figure 2. Reconstructed touch classes under clean (left) and noisy (right) measurement conditions.")
num(s, 5)
note(s, """[The gap, made visible – 2:19 – ~35s]

Before any classification, here's the problem in one picture. On the left,
Figure 2 shows five touch types reconstructed from clean simulated
measurements: you can see exactly where the contact is in each one, and each
class is clearly separable by its pressure intensity and contact area.

On the right, the same five with realistic noise switched on. The clear
distinction between classes has gone — real-world classification is
difficult even to the human eye, let alone a classifier.""")

# ═══════════════════════════════════════════════════════ 6 — result 1
s = slide()
chrome(s)
heading(s, "Identifying and Closing the Sim-to-Real Gap", kicker="Result 1")
rule(s)
pic(s, FIG / "p_collapse.png", 3.12, 2.05, 7.10)
caption(s, 0.85, 5.68, 11.60, "Figure 3. A 1D-CNN trained and evaluated across clean and noisy conditions.")
tb(s, 0.85, 6.08, 11.60, 0.65,
   "Four architectures, three feature representations, same collapse. "
   "It is not a capacity problem — it is a distribution problem.",
   size=16, color=MUTED, line=1.25)
num(s, 6)
note(s, """[Result 1 – 2:54 – ~47s]

To see how serious this gap is, I trained four classifiers — a random
forest, a support vector machine, a multilayer perceptron and a 1D
convolutional neural network — each on 25,000 simulated touches, clean and
noisy in three combinations.

As Figure 3 shows, on clean data the network reaches ninety-four per cent,
close to published results. Trained clean but tested noisy, it collapses to
twenty per cent — as good as guessing. Train and test on noisy data instead,
and it recovers to seventy-six: a fifty-five-point gain. That pattern held
for all four classifiers, so it isn't something a bigger model fixes on its
own.""")

# ═══════════════════════════════════════════════════════ 7 — result 2, main contribution
s = slide()
chrome(s)
heading(s, "Which Noise Components Matter?", kicker="Result 2", title_h=0.90)
rule(s)
pic(s, FIG / "p_inversion.png", 3.12, 2.02, 7.10)
caption(s, 0.85, 5.65, 11.60,
        "Figure 4. Each noise component scored under its own matched protocol versus full deployment corruption.")
tb(s, 0.85, 6.05, 11.60, 0.70,
   "Blue is the conventional protocol, and it is misleading: it rewards a model for "
   "defining an easier world, not for surviving a real one.",
   size=16, color=MUTED, line=1.25)
num(s, 7)
note(s, """[Result 2, the main contribution – 3:41 – ~51s]

The ablation study tests each noise source alone. There are two ways to
score it.

Test each model on the same noise it trained with — the blue bars in Figure
4 — and Gaussian noise looks harmless at ninety-seven per cent, electrode
bias the most damaging.

Or test against the full corruption a real device produces — the orange
bars. The story reverses: electrode bias is the only source anywhere close
to surviving, the rest fall close to chance. That's not a surprise from
nowhere — electrode positioning error is already flagged in the EIT
literature as a dominant modelling error (Kolehmainen et al., 1997). This
result confirms it, under the correct protocol.""")

# ═══════════════════════════════════════════════════════ 8 — result 3
s = slide()
chrome(s)
heading(s, "The Optimal Noise Model", kicker="Result 3")
rule(s)
bullets = (
    "•  Electrode bias + Gaussian noise reaches 76.5% — level with the full "
    "four-component model.\n\n"
    "•  Contact impedance and quantisation add nothing measurable.\n\n"
    "•  The 'no contact' class is an exact zero vector in simulation. Remove it, "
    "and electrode bias alone matches the full model (70.5% vs 69.8%).\n\n"
    "•  Gaussian noise was never conferring robustness — it was making one "
    "degenerate class detectable."
)
tb(s, 0.85, 2.22, 5.55, 4.00, bullets, size=16, line=1.28, space_after=4)
pic(s, FIG / "p_confusion.png", 7.30, 2.05, 4.55)
caption(s, 6.90, 6.02, 5.60,
        "Figure 5. Confusion matrix, noise-aware CNN (Noisy→Deployment, 75.8%).")
tb(s, 0.85, 6.63, 11.60, 0.40,
   "One component to tell touches apart. A noise floor to notice touch at all.",
   size=17.5, bold=True, color=BLUE)
num(s, 8)
note(s, """[Result 3 – 4:32 – ~57s]

Using the correct protocol, two of the four components carry everything:
electrode bias plus Gaussian noise reaches 76.5 per cent, level with the full
model — the noise model halves with no loss of accuracy.

That decomposes further. The 'no contact' class in simulation is an exact
vector of zeros, so I removed it and reran on the four real touches alone.
Electrode bias alone then matches the full model — Gaussian noise was never
providing robustness, just making a class that's zero in training
recognisable when it arrives noisy.

Figure 5 shows why that matters: no-contact is recovered perfectly, a direct
consequence of that Gaussian floor, while point contact is the weakest
class. One component to discriminate between touches, a noise floor to
detect touch at all.""")

# ═══════════════════════════════════════════════════════ 9 — limitations
s = slide()
chrome(s)
heading(s, "What this is, and what it is not", kicker="Limitations", title_h=0.90)
rule(s)
tb(s, 0.85, 2.35, 5.40, 0.35, "What the research supports", size=18, bold=True, color=BLUE)
tb(s, 0.85, 2.90, 5.40, 3.30,
   "Which noise sources a training model needs, and why.\n\n"
   "That the conventional ablation protocol inverts that answer.\n\n"
   "A mechanism for the architecture result, tested rather than assumed.\n\n"
   "An open, reproducible pipeline.",
   size=16, line=1.28, space_after=4)
tb(s, 6.90, 2.35, 5.50, 0.35, "What it doesn't", size=18, bold=True, color=ORANGE)
tb(s, 6.90, 2.90, 5.50, 3.30,
   "Every number here is from simulation. There is no hardware validation.\n\n"
   "Electrode bias is assumed independent per electrode; a stretched sheet would correlate it.\n\n"
   "76% is not deployable. Detecting that contact happened is reliable; identifying which "
   "touch it was is not.",
   size=16, color=MUTED, line=1.28, space_after=4)
num(s, 9)
note(s, """[Limitations – 5:28 – ~54s]

To be upfront about the boundaries: I can claim which noise sources matter
and why, that the conventional protocol inverts that answer, a mechanism for
the architecture result that's tested rather than assumed, and a pipeline
anyone can rerun.

I cannot claim any of it survives contact with hardware. Every number is
simulated in 2D — a real device isn't, and I only classify single, static
frames, so drift over time isn't represented at all. My electrode bias is
drawn independently per electrode, and on a stretched sheet it would be
spatially correlated — measuring that covariance is the first experiment I'd
run next. And seventy-six per cent isn't deployable: detecting contact is
reliable, saying which touch it was is not.""")

# ═══════════════════════════════════════════════════════ 10 — close / conclusions
s = slide()
chrome(s)
heading(s, "Conclusions", kicker="Close")
rule(s)
headline = [
    "Realistic noise collapses simulated accuracy from 94% to 20%; training on it "
    "recovers 76%.",
    "How you evaluate an ablation decides its conclusion, not just its result.",
    "Electrode bias and Gaussian noise, together, are necessary and sufficient.",
]
y = 2.28
for i, line_ in enumerate(headline, start=1):
    tb(s, 0.85, y, 0.45, 0.45, str(i), size=18, bold=True, color=BLUE)
    tb(s, 1.35, y + 0.01, 11.05, 0.62, line_, size=16.5, line=1.22)
    y += 0.72
box(s, 0.85, 4.55, 11.50, 1.55)
tb(s, 1.25, 4.72, 10.80, 1.25,
   "In one sentence: noise-aware training closes most of the simulation-to-reality gap — "
   "but only if you test it against the world the device will actually meet.",
   size=19, bold=True, line=1.28)
tb(s, 0.85, 6.38, 11.30, 0.45, "Thank you. I'm happy to take questions.",
   size=17, color=MUTED)
num(s, 10)
note(s, """[Close – 6:22 – ~31s]

To pull that together: realistic noise collapses accuracy from ninety-four
to twenty per cent, and training on it recovers seventy-six. How you
evaluate an ablation decides its conclusion, not just its result — the
finding I'd defend hardest.

In one sentence: noise-aware training closes most of the gap, but only if
you test it against the world the device will actually meet. Thank you, I'm
happy to take questions.""")

# ═══════════════════════════════════════════════════════ 11 — references
s = slide()
chrome(s)
heading(s, "References", kicker="References")
rule(s)
refs = [
    "Biddiss, E. and Chau, T. (2007) 'Upper-limb prosthetics: critical factors in device "
    "abandonment', American Journal of Physical Medicine & Rehabilitation, 86(12), pp. 977–987.",
    "Chen, H., Langlois, K., Brancart, J., Roels, E., Verstraten, T. and Vanderborght, B. "
    "(2023) 'A novel physical human–robot interface with pressure distribution measurement "
    "based on electrical impedance tomography', IEEE Sensors Journal, 23(18), pp. 21914–21923.",
    "Costa Cornellà, A., Hardman, D., Costi, L., Brancart, J., Van Assche, G. and Iida, F. "
    "(2023) 'Variable sensitivity multimaterial robotic e-skin combining electronic and ionic "
    "conductivity using electrical impedance tomography', Scientific Reports, 13(1), 20004.",
    "Kolehmainen, V., Vauhkonen, M., Karjalainen, P.A. and Kaipio, J.P. (1997) 'Assessment of "
    "errors in static electrical impedance tomography with adjacent and trigonometric current "
    "patterns', Physiological Measurement, 18(4), pp. 289–303.",
    "McDonald, C.L., Westcott-McCoy, S., Weaver, M.R., Haagsma, J. and Kartin, D. (2021) "
    "'Global prevalence of traumatic non-fatal limb amputation', Prosthetics & Orthotics "
    "International, 45(2), pp. 105–114.",
    "Park, K., Yuk, H., Yang, M., Cho, J., Lee, H. and Kim, J. (2022) 'A biomimetic elastomeric "
    "robot skin using electrical impedance and acoustic tomography for tactile sensing', "
    "Science Robotics, 7(67), eabm7187.",
    "Tobin, J., Fong, R., Ray, A., Schneider, J., Zaremba, W. and Abbeel, P. (2017) 'Domain "
    "randomization for transferring deep neural networks from simulation to the real world', "
    "2017 IEEE/RSJ International Conference on Intelligent Robots and Systems (IROS), pp. 23–30.",
]
tb(s, 0.85, 2.20, 5.60, 4.60, "\n".join(refs[:4]), size=12.5, color=INK, line=1.28, space_after=12)
tb(s, 6.90, 2.20, 5.60, 4.60, "\n".join(refs[4:]), size=12.5, color=INK, line=1.28, space_after=12)
num(s, 11)
note(s, """[References – 6:53 – ~10s, not narrated]

Full references for every claim in this talk, in Harvard style, are listed
here for anyone who would like to follow up.""")

out = HERE / "EE52037_Piku_Presentation.pptx"
prs.save(out)
print("saved:", out)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))

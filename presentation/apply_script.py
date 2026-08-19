"""Inject speaker_script.md into the deck's notes pages, so the script has one
source of truth. Run after build_deck.py."""
import re
from pathlib import Path
from pptx import Presentation

HERE = Path(__file__).parent
md = (HERE / "speaker_script.md").read_text(encoding="utf-8")

# each slide section: "## Slide N — label · start · ~Ns"
blocks = re.split(r"\n## Slide (\d+) — ([^\n]*)\n", md)
notes = {}
for i in range(1, len(blocks) - 1, 3):
    n, header, body = int(blocks[i]), blocks[i + 1], blocks[i + 2]
    body = body.split("\n---")[0].strip()
    notes[n] = f"[{header.strip()}]\n\n{body}"

prs = Presentation(HERE / "EE52037_Piku_Presentation.pptx")
total = 0
for i, s in enumerate(prs.slides, 1):
    if i in notes:
        s.notes_slide.notes_text_frame.text = notes[i]
        w = len(notes[i].split()) - len(notes[i].split("\n")[0].split())
        total += w
        print(f"  slide {i:>2}: {w:>4} words")
prs.save(HERE / "EE52037_Piku_Presentation.pptx")
print(f"  TOTAL   : {total:>4} words")
for wpm in (125, 135, 145):
    secs = total / wpm * 60
    print(f"    at {wpm} wpm -> {int(secs//60)}m {secs % 60:02.0f}s")
